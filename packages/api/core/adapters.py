from __future__ import annotations
from typing import Protocol, List, Optional, Dict, Any, BinaryIO
from pydantic import BaseModel, Field

import uuid

class NormalizedSegment(BaseModel):
    """
    Unified representation of a document segment.
    """
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    heading: Optional[str] = None
    heading_chain: List[str] = []
    content: str
    metadata: Dict[str, Any] = {}

class NormalizedDocument(BaseModel):
    """
    The output format for all adapters.
    """
    filename: str
    doc_type: str = "generic"
    segments: List[NormalizedSegment]
    metadata: Dict[str, Any] = {}

class FormatAdapter(Protocol):
    """
    Protocol for document format adapters.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ...

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        ...

# ── Default Text Adapter ──────────────────────────────────────────────────────

class TextAdapter:
    """
    Fallback adapter for plain text and markdown.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        return ext in ["txt", "md"] or mime_type.startswith("text/")

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        raw_text = stream.read().decode("utf-8", errors="replace")
        
        # For now, we use a simple single-segment representation.
        # The complex chunking logic remains in ingest.py for now, 
        # or we could move it here later.
        return NormalizedDocument(
            filename=filename,
            segments=[NormalizedSegment(content=raw_text)]
        )

# ── Excel / CSV Adapter ────────────────────────────────────────────────────────

class ExcelAdapter:
    """
    P4.3-B: Adapter for Excel and CSV files.
    Supports multi-sheet and automatic row/table mode detection.
    B-3: parse_with_config() supports user-specified sheet selection and column mapping.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        return ext in ["xlsx", "xls", "csv"] or "excel" in mime_type or "csv" in mime_type

    def _read_sheets(self, stream: BinaryIO, filename: str) -> dict:
        """Read all sheets from an Excel/CSV file into a dict of DataFrames."""
        import pandas as pd
        import io
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        if ext == "csv":
            content = stream.read().decode("utf-8", errors="replace")
            delimiter = ","
            if "\t" in content and content.count("\t") > content.count(","):
                delimiter = "\t"
            df = pd.read_csv(io.StringIO(content), sep=delimiter)
            if len(df.columns) > 100 or len(df) > 10000:
                raise ValueError("CSV complexity limit exceeded (max 100 columns, 10000 rows)")
            return {"Sheet1": df}
        else:
            dfs = pd.read_excel(stream, sheet_name=None)
            total_cells = 0
            for sn, df in dfs.items():
                if len(df.columns) > 100 or len(df) > 10000:
                    raise ValueError(f"Excel sheet '{sn}' complexity limit exceeded (max 100 columns, 10000 rows)")
                total_cells += len(df.columns) * len(df)
            if total_cells > 500000:
                raise ValueError("Excel total complexity limit exceeded")
            return dfs

    def _process_sheet(self, sheet_name: str, df, title_col: str = None,
                       desc_col: str = None, tag_col: str = None,
                       mode: str = None) -> list[NormalizedSegment]:
        """Process a single sheet into segments using given or auto-detected settings."""
        if df.empty:
            return []

        cols = [str(c) for c in df.columns]
        cols_lower = [c.lower() for c in cols]

        # Auto-detect columns if not specified
        if not title_col:
            for c, cl in zip(cols, cols_lower):
                if any(kw in cl for kw in ["title", "name", "主題", "名稱", "標題"]):
                    title_col = c
                    break
        if not desc_col:
            for c, cl in zip(cols, cols_lower):
                if any(kw in cl for kw in ["description", "desc", "說明", "描述"]):
                    desc_col = c
                    break
        if not tag_col:
            for c, cl in zip(cols, cols_lower):
                if any(kw in cl for kw in ["tag", "label", "標籤", "分類"]):
                    tag_col = c
                    break

        # Auto-detect mode if not specified
        if not mode:
            mode = "row" if len(cols) <= 8 and title_col else "table"

        segments = []
        if mode == "row" and title_col:
            for idx, row in df.iterrows():
                row_data = row.to_dict()
                title_val = str(row_data.get(title_col, f"Row {idx+1}"))

                # Build body from desc_col + other fields
                content_lines = []
                if desc_col and row_data.get(desc_col):
                    content_lines.append(str(row_data[desc_col]))
                    content_lines.append("")

                for k, v in row_data.items():
                    if k in [title_col, desc_col, tag_col]:
                        continue
                    if v is not None and str(v).strip():
                        content_lines.append(f"- **{k}**: {v}")

                # Extract tags from tag column
                tags = []
                if tag_col and row_data.get(tag_col):
                    raw_tags = str(row_data[tag_col])
                    tags = [t.strip() for t in raw_tags.replace("；", ",").replace(";", ",").split(",") if t.strip()]

                segments.append(NormalizedSegment(
                    heading=title_val,
                    heading_chain=[str(sheet_name)],
                    content="\n".join(content_lines),
                    metadata={
                        "sheet": str(sheet_name), "row": int(idx),
                        "tags": tags, "mode": "row"
                    }
                ))
        else:
            # Table mode: Whole sheet as one markdown table
            md_table = df.to_markdown(index=False)
            segments.append(NormalizedSegment(
                heading=str(sheet_name),
                content=md_table,
                metadata={"sheet": str(sheet_name), "mode": "table"}
            ))

        return segments

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        """Default parsing: process all sheets with auto-detected settings."""
        df_dict = self._read_sheets(stream, filename)
        segments = []
        for sheet_name, df in df_dict.items():
            segments.extend(self._process_sheet(str(sheet_name), df))
        return NormalizedDocument(filename=filename, segments=segments)

    async def parse_with_config(self, stream: BinaryIO, filename: str, config: dict) -> NormalizedDocument:
        """
        B-3: Parse with user-specified configuration.
        config schema:
        {
            "selected_sheets": ["Sheet1", "Sheet3"],  // null = all sheets
            "column_mapping": {
                "Sheet1": {
                    "title_col": "Name",
                    "desc_col": "Description",
                    "tag_col": "Category",
                    "mode": "row"  // "row" | "table"
                }
            }
        }
        """
        df_dict = self._read_sheets(stream, filename)
        selected = config.get("selected_sheets")
        mapping = config.get("column_mapping", {})

        segments = []
        for sheet_name, df in df_dict.items():
            sn = str(sheet_name)
            # B-3: Skip unselected sheets
            if selected is not None and sn not in selected:
                continue
            sheet_cfg = mapping.get(sn, {})
            segments.extend(self._process_sheet(
                sn, df,
                title_col=sheet_cfg.get("title_col"),
                desc_col=sheet_cfg.get("desc_col"),
                tag_col=sheet_cfg.get("tag_col"),
                mode=sheet_cfg.get("mode"),
            ))

        return NormalizedDocument(filename=filename, segments=segments)

# ── Word (.docx) Adapter ──────────────────────────────────────────────────────

class WordAdapter:
    """
    P4.3-E1: Adapter for Word documents.
    Preserves heading structure as segments.
    Converts tables to Markdown format.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".docx") or "officedocument.wordprocessingml" in mime_type

    def _table_to_markdown(self, table) -> str:
        """Convert a python-docx Table to a Markdown table string."""
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append(cells)
        if not rows:
            return ""
        # Header
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
        body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        return f"{header}\n{separator}\n{body}"

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        from docx import Document
        from docx.oxml.ns import qn
        doc = Document(stream)
        
        segments = []
        current_heading = "Title"
        current_content = []
        
        # Iterate document body in order (paragraphs and tables interleaved)
        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
            
            if tag == "p":
                # Paragraph element
                from docx.text.paragraph import Paragraph
                para = Paragraph(element, doc)
                if para.style.name.startswith("Heading"):
                    # Save previous segment
                    if current_content:
                        segments.append(NormalizedSegment(
                            heading=current_heading,
                            content="\n".join(current_content)
                        ))
                    current_heading = para.text
                    current_content = []
                else:
                    if para.text.strip():
                        current_content.append(para.text)
            elif tag == "tbl":
                # Table element — convert to Markdown
                from docx.table import Table
                table = Table(element, doc)
                md = self._table_to_markdown(table)
                if md:
                    current_content.append("")
                    current_content.append(md)
                    current_content.append("")
                    
        # Last segment
        if current_content:
            segments.append(NormalizedSegment(
                heading=current_heading,
                content="\n".join(current_content)
            ))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── PowerPoint (.pptx) Adapter ────────────────────────────────────────────────

class PPTAdapter:
    """
    P4.3-E2: Adapter for PowerPoint presentations.
    Each slide becomes a segment.
    Smart merge: consecutive slides with the same title prefix merge into one segment.
    Transition/cover pages (< 30 chars total) are skipped.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".pptx") or "officedocument.presentationml" in mime_type

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        from pptx import Presentation
        prs = Presentation(stream)
        
        # First pass: extract raw slide data
        raw_slides = []
        for i, slide in enumerate(prs.slides, 1):
            title = f"Slide {i}"
            content_parts = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape == slide.shapes.title:
                        title = shape.text
                    else:
                        content_parts.append(shape.text)
            
            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    content_parts.append(f"\nNotes:\n{notes}")
            
            total_text = title + " " + " ".join(content_parts)
            
            # E-2: Skip transition/cover pages with < 30 chars of content
            if len(total_text.strip()) < 30:
                continue
            
            raw_slides.append({
                "title": title,
                "content": "\n".join(content_parts),
                "slide_no": i,
            })

        # Second pass: Smart merge — consecutive slides with same title prefix
        segments = []
        if raw_slides:
            current = raw_slides[0].copy()
            
            def _title_prefix(t: str) -> str:
                """Extract a prefix for grouping (e.g. '需求分析 (1/3)' → '需求分析')."""
                import re
                # Remove trailing numbers, parenthesized ranges, etc.
                clean = re.sub(r'\s*[\(（]\s*\d+\s*/\s*\d+\s*[\)）]\s*$', '', t)
                clean = re.sub(r'\s*[-–—]\s*(續|continued|cont)\s*\.?\s*$', '', clean, flags=re.IGNORECASE)
                return clean.strip()
            
            for slide in raw_slides[1:]:
                curr_prefix = _title_prefix(current["title"])
                next_prefix = _title_prefix(slide["title"])
                
                if curr_prefix and curr_prefix == next_prefix:
                    # Merge: append content to current segment
                    current["content"] += f"\n\n---\n\n{slide['content']}"
                else:
                    # Save current and start new
                    segments.append(NormalizedSegment(
                        heading=current["title"],
                        content=current["content"],
                        metadata={"slide_no": current["slide_no"]}
                    ))
                    current = slide.copy()
            
            # Save last segment
            segments.append(NormalizedSegment(
                heading=current["title"],
                content=current["content"],
                metadata={"slide_no": current["slide_no"]}
            ))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── PDF Adapter ──────────────────────────────────────────────────────────────

class PDFAdapter:
    """
    P4.3-A: Adapter for PDF documents.
    Supports font-based heading detection and OCR fallback.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".pdf") or mime_type == "application/pdf"

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        import pdfplumber
        
        segments = []
        page_count = 0
        has_ocr = False
        
        with pdfplumber.open(stream) as pdf:
            page_count = len(pdf.pages)
            if page_count > 500:
                raise ValueError("PDF too large (max 500 pages)")
            
            current_heading = "Page 1"
            current_content = []
            
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                
                # A-4: OCR Fallback
                if not text or len(text.strip()) < 50:
                    try:
                        import pytesseract
                        from PIL import Image
                        # Convert page to image and OCR
                        img = page.to_image().original
                        text = pytesseract.image_to_string(img, lang="chi_tra+eng")
                        has_ocr = True
                    except Exception as e:
                        print(f"[PDFAdapter] OCR failed for page {i}: {e}")
                        text = f"(OCR Failed for Page {i})"

                # A-3: Font-based Heading reconstruction
                # Use statistical thresholds: mean + 1.5σ → H2, mean + 0.5σ → H3
                words = page.extract_words(extra_attrs=["size", "fontname"])
                lines = []
                if words:
                    # Group words into lines (simple y-axis grouping)
                    current_line = []
                    last_top = words[0]["top"]
                    for w in words:
                        if abs(w["top"] - last_top) > 5:
                            lines.append(current_line)
                            current_line = []
                            last_top = w["top"]
                        current_line.append(w)
                    lines.append(current_line)

                    # Compute font-size statistics across the page
                    all_sizes = [w["size"] for w in words if w.get("size")]
                    if all_sizes:
                        import statistics
                        mean_size = statistics.mean(all_sizes)
                        stdev_size = statistics.pstdev(all_sizes) if len(all_sizes) > 1 else 0
                        h2_threshold = mean_size + 1.5 * stdev_size
                        h3_threshold = mean_size + 0.5 * stdev_size
                    else:
                        h2_threshold = 14
                        h3_threshold = 12

                    for line_words in lines:
                        line_text = " ".join([w["text"] for w in line_words])
                        avg_size = sum([w["size"] for w in line_words]) / len(line_words)
                        
                        # Statistical heading detection
                        is_heading = (
                            avg_size > h3_threshold
                            and stdev_size > 0  # skip if all text is the same size
                            and len(line_text) < 100
                        )
                        if is_heading:
                            if current_content:
                                segments.append(NormalizedSegment(
                                    heading=current_heading,
                                    content="\n".join(current_content),
                                    metadata={"page": i}
                                ))
                            current_heading = line_text
                            current_content = []
                        else:
                            current_content.append(line_text)
                else:
                    # No words, just use the raw text
                    current_content.append(text)

            # Last segment
            if current_content:
                segments.append(NormalizedSegment(
                    heading=current_heading,
                    content="\n".join(current_content),
                    metadata={"page": page_count}
                ))

        return NormalizedDocument(
            filename=filename,
            segments=segments,
            metadata={"page_count": page_count, "has_ocr": has_ocr}
        )

# ── OpenAPI Adapter ─────────────────────────────────────────────────────────

class OpenAPIAdapter:
    """
    P4.3-C1: Adapter for OpenAPI (Swagger) specs.
    Extracts endpoints and schemas as segments, with automatic edge proposals
    between endpoints and the schemas they reference.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        return ext in ["yaml", "yml", "json"] and ("openapi" in filename.lower() or "swagger" in filename.lower())

    def _collect_refs(self, obj: Any) -> list[str]:
        """Recursively collect all $ref schema names from a JSON object."""
        refs = []
        if isinstance(obj, dict):
            if "$ref" in obj:
                ref_path = obj["$ref"]
                # "#/components/schemas/UserProfile" → "UserProfile"
                if "/schemas/" in ref_path:
                    refs.append(ref_path.rsplit("/", 1)[-1])
            for v in obj.values():
                refs.extend(self._collect_refs(v))
        elif isinstance(obj, list):
            for item in obj:
                refs.extend(self._collect_refs(item))
        return refs

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        import yaml
        import json
        
        raw = stream.read().decode("utf-8", errors="replace")
        try:
            spec = json.loads(raw)
        except Exception:
            spec = yaml.safe_load(raw)
            
        segments = []
        if "paths" in spec:
            for path, methods in spec["paths"].items():
                for method, details in methods.items():
                    if method.lower() not in ["get", "post", "put", "delete", "patch"]: continue
                    summary = details.get("summary", details.get("description", ""))

                    # C-1: Resolve $ref edges — endpoint → schema
                    edges = []
                    # Request body refs → "uses" edge
                    req_body = details.get("requestBody", {})
                    for ref_name in self._collect_refs(req_body):
                        edges.append({"to_title_en": f"Schema: {ref_name}", "relation": "depends_on"})

                    # Response refs → "produces" edge
                    responses = details.get("responses", {})
                    for ref_name in self._collect_refs(responses):
                        edges.append({"to_title_en": f"Schema: {ref_name}", "relation": "related_to"})

                    segments.append(NormalizedSegment(
                        heading=f"API: {method.upper()} {path}",
                        content=f"{summary}\n\n```json\n{json.dumps(details, indent=2)}\n```",
                        metadata={"type": "api", "path": path, "method": method, "suggested_edges": edges}
                    ))
                    
        if "components" in spec and "schemas" in spec["components"]:
            for name, schema in spec["components"]["schemas"].items():
                # Schema → Schema refs
                edges = []
                for ref_name in self._collect_refs(schema):
                    if ref_name != name:  # avoid self-ref
                        edges.append({"to_title_en": f"Schema: {ref_name}", "relation": "related_to"})

                segments.append(NormalizedSegment(
                    heading=f"Schema: {name}",
                    content=f"```json\n{json.dumps(schema, indent=2)}\n```",
                    metadata={"type": "schema", "name": name, "suggested_edges": edges}
                ))

        # JSON Schema support ($defs / definitions)
        for defs_key in ["$defs", "definitions"]:
            if defs_key in spec:
                for name, schema in spec[defs_key].items():
                    edges = []
                    for ref_name in self._collect_refs(schema):
                        if ref_name != name:
                            edges.append({"to_title_en": f"Schema: {ref_name}", "relation": "related_to"})
                    segments.append(NormalizedSegment(
                        heading=f"Schema: {name}",
                        content=f"```json\n{json.dumps(schema, indent=2)}\n```",
                        metadata={"type": "schema", "name": name, "suggested_edges": edges}
                    ))

        return NormalizedDocument(filename=filename, segments=segments)

# ── SQL DDL Adapter ──────────────────────────────────────────────────────────

class SQLAdapter:
    """
    P4.3-C2: Adapter for SQL DDL files.
    Each CREATE TABLE becomes a segment. Foreign keys produce 'references' edges.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".sql")

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        raw = stream.read().decode("utf-8", errors="replace")
        
        import re
        table_regex = re.compile(r'CREATE TABLE\s+(?:IF NOT EXISTS\s+)?(\w+)\s*\(([\s\S]*?)\);', re.IGNORECASE)
        # C-2: FK regex — matches REFERENCES target_table(column)
        fk_regex = re.compile(r'REFERENCES\s+(\w+)\s*\(', re.IGNORECASE)
        
        segments = []
        for m in table_regex.finditer(raw):
            table_name = m.group(1)
            columns = m.group(2).strip()

            # C-2: Detect foreign key edges
            edges = []
            for fk_match in fk_regex.finditer(columns):
                ref_table = fk_match.group(1)
                if ref_table.lower() != table_name.lower():
                    edges.append({"to_title_en": f"Table: {ref_table}", "relation": "related_to"})

            segments.append(NormalizedSegment(
                heading=f"Table: {table_name}",
                content=f"Columns:\n```sql\n{columns}\n```",
                metadata={"type": "table", "name": table_name, "suggested_edges": edges}
            ))
            
        if not segments:
            # Fallback to whole file
            segments.append(NormalizedSegment(content=raw))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── Code Adapter ─────────────────────────────────────────────────────────────

class CodeAdapter:
    """
    P4.3-D: Adapter for source code files.
    Extracts structure (classes, functions) as segments.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        return ext in ["py", "js", "ts", "tsx", "jsx", "go", "java", "cpp", "h"]

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        raw = stream.read().decode("utf-8", errors="replace")
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        
        segments = []
        
        # D-3: Collect all symbol names for basic dependency detection
        symbol_names = set()
        
        if ext == "py":
            import ast
            try:
                tree = ast.parse(raw)
                for node in tree.body:
                    if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                        symbol_names.add(node.name)
            except: pass
        elif ext in ["js", "ts", "tsx", "jsx"]:
            import re
            # Match: function name, class name, const name = () =>
            js_sym_regex = re.compile(r'(?:function|class|const|let|var)\s+([a-zA-Z0-9_$]+)', re.MULTILINE)
            for m in js_sym_regex.finditer(raw):
                symbol_names.add(m.group(1))

        if ext == "py":
            import ast
            try:
                tree = ast.parse(raw)
                lines = raw.splitlines()
                for node in tree.body:
                    if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                        name = node.name
                        type_name = "Class" if isinstance(node, ast.ClassDef) else "Function"
                        start_line = node.lineno
                        end_line = getattr(node, "end_lineno", start_line + 5)
                        node_source = "\n".join(lines[start_line-1 : end_line])
                        
                        # D-3: Simple dependency detection (mentions)
                        deps = [s for s in symbol_names if s != name and s in node_source]
                        edges = [{"to_title_en": f"Function: {d}", "relation": "depends_on"} for d in deps]
                        
                        segments.append(NormalizedSegment(
                            heading=f"{type_name}: {name}",
                            content=f"```python\n{node_source}\n```",
                            metadata={"type": type_name.lower(), "name": name, "line": start_line, "suggested_edges": edges}
                        ))
            except Exception as e:
                print(f"[CodeAdapter] Python parse failed: {e}")
                segments.append(NormalizedSegment(content=raw))
                
        elif ext in ["js", "ts", "tsx", "jsx"]:
            # D-2: JS/TS Extraction using TypeScript compiler API
            import subprocess
            import json
            import tempfile
            import os
            
            # Write content to a proper temp file for the parser
            try:
                with tempfile.NamedTemporaryFile(mode="w", suffix=f".{ext}", encoding="utf-8", delete=False) as tmp:
                    tmp.write(raw)
                    temp_path = tmp.name
                
                # Call the upgraded TS parser
                result = subprocess.run(
                    ["node", "core/ts_parser.js", temp_path],
                    capture_output=True, text=True, check=False
                )
                
                if result.returncode == 0:
                    symbols = json.loads(result.stdout)
                    lines = raw.splitlines()

                    # Map kind to heading prefix
                    _kind_label = {
                        "class": "Class", "interface": "Interface", "type": "Type",
                        "enum": "Enum", "function": "Function", "method": "Method",
                        "component": "Component", "variable": "Symbol", "export": "Export",
                    }

                    for sym in symbols:
                        name = sym["name"]
                        kind = sym.get("kind", "symbol")
                        line = sym["line"]
                        end_line = sym.get("endLine", line + 10)
                        
                        # Use accurate line range from parser
                        start = max(0, line - 1)
                        end = min(len(lines), end_line)
                        content_block = "\n".join(lines[start:end])
                        
                        # D-3: Dependency detection — string matching
                        base_name = name.split(".")[-1] if "." in name else name
                        deps = [s for s in symbol_names if s != base_name and s in content_block]
                        edges = [{"to_title_en": f"Symbol: {d}", "relation": "depends_on"} for d in deps]

                        # D-3: Extends edges from AST heritage
                        extends_list = sym.get("extends", [])
                        for parent_name in extends_list:
                            edges.append({"to_title_en": f"Class: {parent_name}", "relation": "extends"})

                        label = _kind_label.get(kind, "Symbol")
                        heading = f"{label}: {name}"
                        doc = sym.get("doc", "")
                        doc_line = f"\n\n{doc}" if doc else ""

                        segments.append(NormalizedSegment(
                            heading=heading,
                            content=f"{doc_line}\n```{ext}\n{content_block}\n```".strip(),
                            metadata={
                                "name": name, "kind": kind,
                                "line": line, "endLine": end_line,
                                "parent": sym.get("parent"),
                                "suggested_edges": edges
                            }
                        ))
                else:
                    print(f"[CodeAdapter] JS Parser error: {result.stderr}")
                    segments.append(NormalizedSegment(content=raw))
            except Exception as e:
                print(f"[CodeAdapter] JS Subprocess failed: {e}")
                segments.append(NormalizedSegment(content=raw))
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        else:
            segments.append(NormalizedSegment(content=raw))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── GitHub Adapter ────────────────────────────────────────────────────────────

class GitHubAdapter:
    """
    P4.3-G2: Adapter for GitHub URLs.
    Fetches PR or Issue content and comments.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        # Check if the file is a .url or .txt containing a github link
        return filename.endswith(".github") or "github.com" in filename.lower()

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        import httpx
        url = stream.read().decode("utf-8").strip()
        if "github.com" not in url:
            return NormalizedDocument(filename=filename, segments=[NormalizedSegment(content=url)])
            
        # Basic URL parsing
        # github.com/owner/repo/pull/123 -> api.github.com/repos/owner/repo/pulls/123
        api_url = url.replace("github.com", "api.github.com/repos")
        if "/pull/" in api_url:
            api_url = api_url.replace("/pull/", "/pulls/")
        elif "/issues/" in api_url:
            pass # already correct
            
        headers = {"Accept": "application/vnd.github.v3+json"}
        # Note: In production, we'd use a GitHub App token or user token from DB
        
        async with httpx.AsyncClient() as client:
            resp = await client.get(api_url, headers=headers)
            if resp.status_code != 200:
                return NormalizedDocument(filename=filename, segments=[NormalizedSegment(content=f"Error fetching GitHub: {resp.text}")])
            
            data = resp.json()
            title = data.get("title", "GitHub Item")
            body = data.get("body", "")
            
            segments = [NormalizedSegment(
                heading=f"GitHub: {title}",
                content=f"{body}\n\nURL: {url}",
                metadata={"source": "github", "url": url}
            )]
            
            # Optionally fetch comments
            comments_url = data.get("comments_url")
            if comments_url:
                c_resp = await client.get(comments_url, headers=headers)
                if c_resp.status_code == 200:
                    comments = c_resp.json()
                    for i, c in enumerate(comments, 1):
                        segments.append(NormalizedSegment(
                            heading=f"Comment {i} by {c['user']['login']}",
                            content=c['body'],
                            metadata={"source": "github", "type": "comment"}
                        ))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── Git / Diff Adapter ────────────────────────────────────────────────────────

class GitAdapter:
    """
    P4.3-G: Adapter for Git diff/patch files.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".diff") or filename.endswith(".patch")

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        raw = stream.read().decode("utf-8", errors="replace")
        
        # Split by file diffs
        import re
        file_diff_regex = re.compile(r'^diff --git a/.* b/(.*)$', re.MULTILINE)
        
        segments = []
        last_pos = 0
        current_file = "Unknown File"
        
        for m in file_diff_regex.finditer(raw):
            if last_pos > 0:
                segments.append(NormalizedSegment(
                    heading=f"Diff: {current_file}",
                    content=f"```diff\n{raw[last_pos:m.start()]}\n```",
                    metadata={"file": current_file}
                ))
            current_file = m.group(1)
            last_pos = m.start()
            
        # Last segment
        if last_pos < len(raw):
            segments.append(NormalizedSegment(
                heading=f"Diff: {current_file}",
                content=f"```diff\n{raw[last_pos:]}\n```",
                metadata={"file": current_file}
            ))
            
        return NormalizedDocument(filename=filename, segments=segments)

# ── Slack Export Adapter ──────────────────────────────────────────────────────

# F-2: Signal detection keywords for decisions / action items
_SIGNAL_PATTERNS = {
    "decision": [
        r"(?:決定|確認|agreed|decided|conclusion|結論|we\'ll go with|let\'s go with|final(?:ly|ized)?)",
    ],
    "action_item": [
        r"(?:TODO|action item|負責|請.*處理|@\w+\s+(?:請|will|should|needs? to)|follow[- ]?up)",
    ],
    "confirmation": [
        r"(?:confirmed?|approved?|LGTM|looks good|沒問題|OK\b|可以|同意)",
    ],
}


class SlackAdapter:
    """
    P4.3-F: Adapter for Slack export ZIPs.
    Implements two-stage extraction:
      - Stage 1 (AI summary): deferred to process_ingestion (requires provider).
      - Stage 2 (signal detection): regex keyword scan done here.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.endswith(".zip") and ("slack" in filename.lower() or "export" in filename.lower())

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        import zipfile
        import json
        import re

        segments = []
        with zipfile.ZipFile(stream) as z:
            total_uncompressed = 0
            total_compressed = sum(info.compress_size for info in z.infolist()) or 1
            
            for info in z.infolist():
                if ".." in info.filename or info.filename.startswith("/") or info.filename.startswith("\\"):
                    raise ValueError("Invalid path in ZIP (Path Traversal attempt)")
                total_uncompressed += info.file_size

            if total_uncompressed / total_compressed > 100:
                raise ValueError("Zip compression ratio too high (Zip Bomb detected)")
                
            if len(z.infolist()) > 5000 or total_uncompressed > 500 * 1024 * 1024:
                raise ValueError("Zip contents too large or too many entries")
                
            for name in z.namelist():
                if name.endswith(".json") and not name.startswith("metadata"):
                    try:
                        with z.open(name) as f:
                            messages = json.load(f)

                        if not isinstance(messages, list) or not messages:
                            continue

                        # ── Group messages by thread_ts ────────────────────
                        threads: Dict[str, list] = {}
                        for msg in messages:
                            ts = msg.get("thread_ts") or msg.get("ts") or "root"
                            threads.setdefault(ts, []).append(msg)

                        for thread_ts, thread_msgs in threads.items():
                            lines = []
                            participants = set()
                            timestamps = []
                            for msg in thread_msgs:
                                user = msg.get("user", msg.get("username", "Unknown"))
                                text = msg.get("text", "")
                                if text:
                                    lines.append(f"[{user}]: {text}")
                                    participants.add(user)
                                ts_val = msg.get("ts")
                                if ts_val:
                                    timestamps.append(ts_val)

                            if not lines:
                                continue

                            full_text = "\n".join(lines)

                            # ── F-2 Stage 2: Signal Detection ─────────────
                            signals: Dict[str, List[str]] = {
                                "decision": [], "action_item": [], "confirmation": []
                            }
                            for category, patterns in _SIGNAL_PATTERNS.items():
                                combined = "|".join(patterns)
                                for line in lines:
                                    if re.search(combined, line, re.IGNORECASE):
                                        signals[category].append(line.strip())

                            signal_section = ""
                            for cat, hits in signals.items():
                                if hits:
                                    label = {"decision": "Decisions", "action_item": "Action Items", "confirmation": "Confirmations"}[cat]
                                    signal_section += f"\n\n### Detected {label}\n" + "\n".join(f"- {h}" for h in hits[:10])

                            # Determine a suitable heading from the first message or topic
                            first_text = thread_msgs[0].get("text", "")[:80]
                            heading = f"Slack Thread: {first_text}" if first_text else f"Slack Thread ({name})"

                            # Build segment with AI-ready content
                            content = full_text
                            if signal_section:
                                content += f"\n\n---{signal_section}"

                            # Hint for the AI extraction pipeline (Stage 1 deferred)
                            if len(lines) > 3:
                                content += (
                                    "\n\n[Extraction Guidance: Summarize this thread. "
                                    "Extract decisions as factual nodes, action items as procedural nodes, "
                                    "and unresolved questions as context nodes.]"
                                )

                            segments.append(NormalizedSegment(
                                heading=heading,
                                content=content,
                                metadata={
                                    "channel": name,
                                    "thread_ts": thread_ts,
                                    "msg_count": len(lines),
                                    "participants": sorted(participants),
                                    "date": timestamps[0] if timestamps else None,
                                    "signals": {k: len(v) for k, v in signals.items()},
                                }
                            ))
                    except Exception:
                        continue

        return NormalizedDocument(filename=filename, segments=segments)

# ── Image Adapter (Screenshots) ─────────────────────────────────────────────

class ImageAdapter:
    """
    P4.3-H: Adapter for images (screenshots, architecture diagrams).
    Uses multimodal LLM to describe the image.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        return filename.split('.')[-1].lower() in ["png", "jpg", "jpeg", "webp"] or mime_type.startswith("image/")

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        from PIL import Image
        import io
        import base64
        
        raw_bytes = stream.read()
        try:
            with Image.open(io.BytesIO(raw_bytes)) as img:
                if img.width > 8000 or img.height > 8000:
                    raise ValueError(f"Image dimensions too large ({img.width}x{img.height}). Max 8000x8000.")
        except Exception as e:
            if isinstance(e, ValueError): raise e
            pass
            
        image_data = base64.b64encode(raw_bytes).decode("utf-8")
        
        # We store the base64 or a reference. For extraction, we'll need to send this to LLM.
        # For now, we return a segment that describes what to do.
        return NormalizedDocument(
            filename=filename,
            segments=[NormalizedSegment(
                heading=f"Image: {filename}",
                content=f"[Attached Image: {filename}]\n(Multimodal analysis will be performed during extraction)",
                metadata={"image_base64": image_data, "mime_type": "image/png"} # simplified
            )]
        )

# ── Email (.eml) Adapter ──────────────────────────────────────────────────────

class EmailAdapter:
    """
    P4.3-F: Adapter for .eml email files.
    Extracts subject, sender, recipients, date, and body text.
    Applies signal detection (reuses SlackAdapter patterns) for decision/action item extraction.
    """
    def can_handle(self, filename: str, mime_type: str) -> bool:
        ext = filename.split('.')[-1].lower() if '.' in filename else ""
        return ext == "eml" or mime_type in ["message/rfc822"]

    async def parse(self, stream: BinaryIO, filename: str) -> NormalizedDocument:
        import email
        from email import policy
        import re

        raw_bytes = stream.read()
        msg = email.message_from_bytes(raw_bytes, policy=policy.default)

        subject = msg.get("Subject", "(No Subject)")
        from_addr = msg.get("From", "")
        to_addr = msg.get("To", "")
        cc_addr = msg.get("Cc", "")
        date_str = msg.get("Date", "")

        # Extract plain text body
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        body += payload.decode(charset, errors="replace")
                elif ctype == "text/html" and not body:
                    # Fallback: strip HTML tags if no plain text part
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        html = payload.decode(charset, errors="replace")
                        body = re.sub(r'<[^>]+>', '', html)
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or "utf-8"
                body = payload.decode(charset, errors="replace")

        body = body.strip()

        # Build header section
        header_lines = [
            f"**From:** {from_addr}",
            f"**To:** {to_addr}",
        ]
        if cc_addr:
            header_lines.append(f"**CC:** {cc_addr}")
        if date_str:
            header_lines.append(f"**Date:** {date_str}")

        content = "\n".join(header_lines) + "\n\n---\n\n" + body

        # Signal detection (reuse Slack patterns)
        signals: Dict[str, list] = {"decision": [], "action_item": [], "confirmation": []}
        for category, patterns in _SIGNAL_PATTERNS.items():
            combined = "|".join(patterns)
            for line in body.splitlines():
                if re.search(combined, line, re.IGNORECASE):
                    signals[category].append(line.strip())

        signal_section = ""
        for cat, hits in signals.items():
            if hits:
                label = {"decision": "Decisions", "action_item": "Action Items", "confirmation": "Confirmations"}[cat]
                signal_section += f"\n\n### Detected {label}\n" + "\n".join(f"- {h}" for h in hits[:10])

        if signal_section:
            content += f"\n\n---{signal_section}"

        # Add extraction guidance for AI pipeline
        if len(body) > 100:
            content += (
                "\n\n[Extraction Guidance: This is an email thread. "
                "Extract key decisions as factual nodes, action items as procedural nodes, "
                "and important context/background as context nodes.]"
            )

        metadata = {
            "format": "email",
            "subject": subject,
            "from": from_addr,
            "to": to_addr,
            "cc": cc_addr,
            "date": date_str,
        }

        return NormalizedDocument(
            filename=filename,
            segments=[NormalizedSegment(
                heading=f"Email: {subject}",
                content=content,
                metadata=metadata,
            )],
            metadata=metadata,
        )

# ── Format Detector ───────────────────────────────────────────────────────────

_ADAPTERS: List[FormatAdapter] = [
    OpenAPIAdapter(),
    SQLAdapter(),
    CodeAdapter(),
    ExcelAdapter(),
    WordAdapter(),
    PPTAdapter(),
    PDFAdapter(),
    GitAdapter(),
    GitHubAdapter(),
    SlackAdapter(),
    EmailAdapter(),
    ImageAdapter(),
    TextAdapter(),
    # Future adapters will be added here
]

def get_adapter_for_file(filename: str, mime_type: str = "application/octet-stream") -> FormatAdapter:
    """
    Core-2: Dispatch the appropriate adapter based on filename or MIME type.
    """
    for adapter in _ADAPTERS:
        if adapter.can_handle(filename, mime_type):
            return adapter
    return TextAdapter() # Fallback
